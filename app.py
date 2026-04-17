# filepath: app.py
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
import csv
import json
import threading
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from deep_translator import GoogleTranslator, DeeplTranslator, LibreTranslator
import sys

# Application info
APP_NAME = "PPT Translator"
VERSION = "1.0.0"

# Default glossary file path
GLOSSARY_FILE = "glossary.json"

class PPTTranslatorApp:
    def __init__(self, root):
        self.root = root
        self.root.title(f"{APP_NAME} v{VERSION}")
        self.root.geometry("800x650")
        self.root.minsize(600, 500)
        
        # Data
        self.selected_file = None
        self.glossary = {}
        
        # Translation settings
        self.translator_service = tk.StringVar(value="google")
        self.custom_endpoint = tk.StringVar(value="")
        self.api_key = tk.StringVar(value="")
        
        # Load glossary
        self.load_glossary()
        
        # Setup UI
        self.setup_ui()
        
    def setup_ui(self):
        """Setup the user interface"""
        # Main container
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Header
        header_frame = ttk.Frame(main_frame)
        header_frame.pack(fill=tk.X, pady=(0, 10))
        
        title_label = ttk.Label(header_frame, text=f"{APP_NAME} - Japanese to English Translator", 
                                font=("Segoe UI", 16, "bold"))
        title_label.pack(side=tk.LEFT)
        
        # File Selection Section
        file_frame = ttk.LabelFrame(main_frame, text="1. Select PowerPoint File", padding="10")
        file_frame.pack(fill=tk.X, pady=5)
        
        self.select_btn = ttk.Button(file_frame, text="Select .pptx File", command=self.select_file)
        self.select_btn.pack(side=tk.LEFT, padx=5)
        
        self.file_label = ttk.Label(file_frame, text="No file selected", foreground="gray")
        self.file_label.pack(side=tk.LEFT, padx=10)
        
        self.slide_count_label = ttk.Label(file_frame, text="", foreground="gray")
        self.slide_count_label.pack(side=tk.LEFT, padx=10)
        
        # Translation Section
        trans_frame = ttk.LabelFrame(main_frame, text="2. Translate", padding="10")
        trans_frame.pack(fill=tk.X, pady=5)
        
        # Service selection
        service_frame = ttk.Frame(trans_frame)
        service_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(service_frame, text="Service:").pack(side=tk.LEFT)
        
        service_combo = ttk.Combobox(service_frame, textvariable=self.translator_service, 
                                      values=["google", "deepl", "libre", "custom"], 
                                      state="readonly", width=12)
        service_combo.pack(side=tk.LEFT, padx=5)
        service_combo.bind("<<ComboboxSelected>>", self.on_service_change)
        
        ttk.Label(service_frame, text="Endpoint:").pack(side=tk.LEFT, padx=(20, 5))
        self.endpoint_entry = ttk.Entry(service_frame, textvariable=self.custom_endpoint, width=25)
        self.endpoint_entry.pack(side=tk.LEFT, padx=5)
        
        ttk.Label(service_frame, text="API Key:").pack(side=tk.LEFT, padx=(10, 5))
        self.apikey_entry = ttk.Entry(service_frame, textvariable=self.api_key, width=20, show="*")
        self.apikey_entry.pack(side=tk.LEFT, padx=5)
        
        # Initial state - disable endpoint and API key fields for Google
        self.endpoint_entry.config(state=tk.DISABLED)
        self.apikey_entry.config(state=tk.DISABLED)
        
        # Translate button and progress
        btn_frame = ttk.Frame(trans_frame)
        btn_frame.pack(fill=tk.X, pady=5)
        
        self.translate_btn = ttk.Button(btn_frame, text="Start Translation", 
                                        command=self.start_translation, state=tk.DISABLED)
        self.translate_btn.pack(side=tk.LEFT, padx=5)
        
        self.progress = ttk.Progressbar(btn_frame, mode='indeterminate', length=200)
        self.progress.pack(side=tk.LEFT, padx=10, fill=tk.X, expand=True)
        
        self.status_label = ttk.Label(btn_frame, text="Ready", foreground="gray")
        self.status_label.pack(side=tk.LEFT, padx=10)
        
        # Glossary Section
        glossary_frame = ttk.LabelFrame(main_frame, text="3. Glossary (Vocabulary)", padding="10")
        glossary_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # Glossary input
        input_frame = ttk.Frame(glossary_frame)
        input_frame.pack(fill=tk.X, pady=5)
        
        ttk.Label(input_frame, text="Japanese:").pack(side=tk.LEFT)
        self.jp_entry = ttk.Entry(input_frame, width=20)
        self.jp_entry.pack(side=tk.LEFT, padx=5)
        
        ttk.Label(input_frame, text="English:").pack(side=tk.LEFT, padx=10)
        self.en_entry = ttk.Entry(input_frame, width=20)
        self.en_entry.pack(side=tk.LEFT, padx=5)
        
        self.add_btn = ttk.Button(input_frame, text="Add", command=self.add_glossary)
        self.add_btn.pack(side=tk.LEFT, padx=10)
        
        # Glossary buttons
        btn_frame = ttk.Frame(glossary_frame)
        btn_frame.pack(fill=tk.X, pady=5)
        
        ttk.Button(btn_frame, text="Import CSV", command=self.import_csv).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="Export CSV", command=self.export_csv).pack(side=tk.LEFT, padx=5)
        ttk.Button(btn_frame, text="Clear All", command=self.clear_glossary).pack(side=tk.LEFT, padx=5)
        
        # Glossary table
        table_frame = ttk.Frame(glossary_frame)
        table_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        
        # Treeview for glossary
        columns = ("Japanese", "English")
        self.glossary_tree = ttk.Treeview(table_frame, columns=columns, show="headings", height=8)
        self.glossary_tree.heading("Japanese", text="Japanese")
        self.glossary_tree.heading("English", text="English")
        self.glossary_tree.column("Japanese", width=200)
        self.glossary_tree.column("English", width=200)
        
        scrollbar = ttk.Scrollbar(table_frame, orient=tk.VERTICAL, command=self.glossary_tree.yview)
        self.glossary_tree.configure(yscrollcommand=scrollbar.set)
        
        self.glossary_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Delete button
        self.delete_btn = ttk.Button(glossary_frame, text="Delete Selected", command=self.delete_glossary)
        self.delete_btn.pack(pady=5)
        
        # Footer
        footer_frame = ttk.Frame(main_frame)
        footer_frame.pack(fill=tk.X, pady=(10, 0))
        
        self.footer_label = ttk.Label(footer_frame, text="", foreground="gray", font=("Segoe UI", 9))
        self.footer_label.pack(side=tk.LEFT)
        
        # Update glossary display
        self.update_glossary_display()
        
    def select_file(self):
        """Select a PowerPoint file"""
        filename = filedialog.askopenfilename(
            title="Select PowerPoint File",
            filetypes=[("PowerPoint Files", "*.pptx"), ("All Files", "*.*")]
        )
        
        if filename:
            self.selected_file = filename
            self.file_label.config(text=os.path.basename(filename), foreground="black")
            
            # Count slides
            try:
                prs = Presentation(filename)
                slide_count = len(prs.slides)
                self.slide_count_label.config(text=f"({slide_count} slides)")
                self.translate_btn.config(state=tk.NORMAL)
                self.status_label.config(text="File loaded")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to read PowerPoint file: {str(e)}")
                self.selected_file = None
                self.translate_btn.config(state=tk.DISABLED)
    
    def load_glossary(self):
        """Load glossary from file"""
        if os.path.exists(GLOSSARY_FILE):
            try:
                with open(GLOSSARY_FILE, 'r', encoding='utf-8') as f:
                    self.glossary = json.load(f)
            except:
                self.glossary = {}
    
    def save_glossary(self):
        """Save glossary to file"""
        try:
            with open(GLOSSARY_FILE, 'w', encoding='utf-8') as f:
                json.dump(self.glossary, f, ensure_ascii=False, indent=2)
        except Exception as e:
            messagebox.showwarning("Warning", f"Failed to save glossary: {str(e)}")
    
    def add_glossary(self):
        """Add a glossary term"""
        jp = self.jp_entry.get().strip()
        en = self.en_entry.get().strip()
        
        if jp and en:
            self.glossary[jp] = en
            self.jp_entry.delete(0, tk.END)
            self.en_entry.delete(0, tk.END)
            self.update_glossary_display()
            self.save_glossary()
    
    def delete_glossary(self):
        """Delete selected glossary term"""
        selected = self.glossary_tree.selection()
        if selected:
            item = self.glossary_tree.item(selected[0])
            jp = item['values'][0]
            if jp in self.glossary:
                del self.glossary[jp]
                self.update_glossary_display()
                self.save_glossary()
    
    def clear_glossary(self):
        """Clear all glossary terms"""
        if messagebox.askyesno("Confirm", "Clear all glossary terms?"):
            self.glossary = {}
            self.update_glossary_display()
            self.save_glossary()
    
    def update_glossary_display(self):
        """Update the glossary table"""
        for item in self.glossary_tree.get_children():
            self.glossary_tree.delete(item)
        
        for jp, en in sorted(self.glossary.items()):
            self.glossary_tree.insert("", tk.END, values=(jp, en))
        
        self.footer_label.config(text=f"Glossary: {len(self.glossary)} terms")
    
    def import_csv(self):
        """Import glossary from CSV"""
        filename = filedialog.askopenfilename(
            title="Import Glossary",
            filetypes=[("CSV Files", "*.csv"), ("All Files", "*.*")]
        )
        
        if filename:
            try:
                with open(filename, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    count = 0
                    for row in reader:
                        if len(row) >= 2:
                            self.glossary[row[0].strip()] = row[1].strip()
                            count += 1
                    self.update_glossary_display()
                    self.save_glossary()
                    messagebox.showinfo("Success", f"Imported {count} terms")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to import: {str(e)}")
    
    def export_csv(self):
        """Export glossary to CSV"""
        filename = filedialog.asksaveasfilename(
            title="Export Glossary",
            defaultextension=".csv",
            filetypes=[("CSV Files", "*.csv"), ("All Files", "*.*")]
        )
        
        if filename:
            try:
                with open(filename, 'w', encoding='utf-8', newline='') as f:
                    writer = csv.writer(f)
                    for jp, en in sorted(self.glossary.items()):
                        writer.writerow([jp, en])
                    messagebox.showinfo("Success", f"Exported {len(self.glossary)} terms")
            except Exception as e:
                messagebox.showerror("Error", f"Failed to export: {str(e)}")
    
    def apply_pre_glossary(self, text):
        """Apply glossary before translation"""
        result = text
        for jp, en in self.glossary.items():
            result = result.replace(jp, en)
        return result
    
    def on_service_change(self, event=None):
        """Handle service selection change"""
        service = self.translator_service.get()
        if service == "custom":
            self.endpoint_entry.config(state=tk.NORMAL)
            self.apikey_entry.config(state=tk.NORMAL)
        elif service == "deepl":
            self.endpoint_entry.config(state=tk.DISABLED)
            self.apikey_entry.config(state=tk.NORMAL)
        elif service == "libre":
            self.endpoint_entry.config(state=tk.DISABLED)
            self.apikey_entry.config(state=tk.DISABLED)
        else:  # google
            self.endpoint_entry.config(state=tk.DISABLED)
            self.apikey_entry.config(state=tk.DISABLED)
    
    def apply_post_glossary(self, text):
        """Apply glossary after translation (for terms that should stay as-is)"""
        # This can be used if user wants to preserve certain English terms
        return text
    
    def translate_text(self, text):
        """Translate text from Japanese to English"""
        if not text.strip():
            return text
        
        # Apply pre-glossary
        text = self.apply_pre_glossary(text)
        
        try:
            service = self.translator_service.get()
            
            if service == "google":
                translator = GoogleTranslator(source='auto', target='en')
                translated = translator.translate(text)
            
            elif service == "deepl":
                api_key = self.api_key.get().strip()
                if api_key:
                    translator = DeeplTranslator(api_key=api_key, source='auto', target='en')
                else:
                    translator = DeeplTranslator(source='auto', target='en')
                translated = translator.translate(text)
            
            elif service == "libre":
                translator = LibreTranslator(source='auto', target='en')
                translated = translator.translate(text)
            
            elif service == "custom":
                endpoint = self.custom_endpoint.get().strip()
                api_key = self.api_key.get().strip()
                
                if not endpoint:
                    raise ValueError("Custom endpoint is required")
                
                # Custom API call
                headers = {"Content-Type": "application/json"}
                if api_key:
                    headers["Authorization"] = f"Bearer {api_key}"
                
                # Try common API formats
                payload = {"text": text, "source": "auto", "target": "en"}
                
                response = requests.post(endpoint, json=payload, headers=headers, timeout=30)
                response.raise_for_status()
                
                # Try to parse response (common formats)
                result = response.json()
                if isinstance(result, dict):
                    translated = result.get("translatedText") or result.get("translation") or result.get("text", text)
                else:
                    translated = str(result)
            else:
                translated = text
            
            # Apply post-glossary
            translated = self.apply_post_glossary(translated)
            
            return translated
        except Exception as e:
            print(f"Translation error: {e}")
            return text
    
    def translate_presentation(self):
        """Translate the PowerPoint presentation"""
        if not self.selected_file:
            return
        
        try:
            self.status_label.config(text="Loading PowerPoint...")
            self.progress.start()
            
            # Load presentation
            prs = Presentation(self.selected_file)
            total_slides = len(prs.slides)
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                self.status_label.config(text=f"Translating slide {slide_idx + 1}/{total_slides}...")
                self.root.update()
                
                # Process all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        original_text = shape.text
                        translated_text = self.translate_text(original_text)
                        shape.text = translated_text
                    
                    # Process tables
                    if shape.has_table:
                        for cell in shape.table.iter_cells():
                            if cell.text:
                                original_text = cell.text
                                translated_text = self.translate_text(original_text)
                                cell.text = translated_text
            
            # Save output file
            output_file = self.selected_file.replace('.pptx', '_translated.pptx')
            counter = 1
            while os.path.exists(output_file):
                output_file = self.selected_file.replace('.pptx', f'_translated_{counter}.pptx')
                counter += 1
            
            self.status_label.config(text="Saving file...")
            prs.save(output_file)
            
            self.progress.stop()
            self.status_label.config(text="Complete!")
            messagebox.showinfo("Success", f"Translation complete!\nSaved to: {os.path.basename(output_file)}")
            
        except Exception as e:
            self.progress.stop()
            self.status_label.config(text="Error")
            messagebox.showerror("Error", f"Translation failed: {str(e)}")
    
    def start_translation(self):
        """Start translation in a separate thread"""
        if not self.selected_file:
            messagebox.showwarning("Warning", "Please select a PowerPoint file first")
            return
        
        # Disable button during translation
        self.translate_btn.config(state=tk.DISABLED)
        self.select_btn.config(state=tk.DISABLED)
        
        # Run in thread
        thread = threading.Thread(target=self.translate_presentation)
        thread.start()
        
        # Check thread completion
        self.root.after(100, self.check_thread, thread)
    
    def check_thread(self, thread):
        """Check if translation thread is complete"""
        if thread.is_alive():
            self.root.after(100, self.check_thread, thread)
        else:
            self.translate_btn.config(state=tk.NORMAL)
            self.select_btn.config(state=tk.NORMAL)
            self.progress.stop()

def main():
    root = tk.Tk()
    app = PPTTranslatorApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()